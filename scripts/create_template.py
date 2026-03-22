"""Create conformant pptgen .pptx template files.

Reads every entry in templates/registry.yaml and generates the .pptx file at
the path each entry declares.  Adding a new template to the registry is the
only change required — this script needs no edits.

Source
------
The branded base is ``template/HC_Powerpoint_Template_with_pptgen_placeholders.potx``.
It already contains all six canonical layouts with the correct pptgen shape names.
This script:
  1. Converts the .potx to a loadable .pptx (content-type fix in the zip).
  2. Strips the three sample slides so the output is a blank-slide base.
  3. Saves one copy per registered template path.

No layout renaming or shape injection is needed — the source template already
follows the Template Authoring Standard.

Run from the repository root:
    python scripts/create_template.py
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

from pptx import Presentation


_REPO_ROOT = Path(__file__).parent.parent
_SOURCE_POTX = _REPO_ROOT / "template" / "HC_Powerpoint_Template_with_pptgen_placeholders.potx"
_REGISTRY_PATH = _REPO_ROOT / "templates" / "registry.yaml"

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _potx_to_pptx_bytes(potx_path: Path) -> bytes:
    """Convert a .potx file to a loadable .pptx by fixing the content type."""
    with open(potx_path, "rb") as f:
        data = f.read()

    buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(data), "r") as zin:
        with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                content = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    content = content.replace(
                        b"presentationml.template.main+xml",
                        b"presentationml.presentation.main+xml",
                    )
                zout.writestr(item, content)
    return buf.getvalue()


def _strip_slides(prs: Presentation) -> None:
    """Remove all existing slides from *prs* in place."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.attrib.get(f"{{{_R_NS}}}id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)


def create_template(output_path: Path, pptx_bytes: bytes) -> None:
    """Save a slide-stripped copy of the branded template to *output_path*."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    _strip_slides(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Created: {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    import sys
    import yaml

    if not _SOURCE_POTX.exists():
        raise SystemExit(f"Source template not found: {_SOURCE_POTX}")
    if not _REGISTRY_PATH.exists():
        raise SystemExit(f"Registry not found: {_REGISTRY_PATH}")

    with open(_REGISTRY_PATH) as f:
        registry_data = yaml.safe_load(f)

    entries = registry_data.get("templates", [])
    if not entries:
        raise SystemExit("No templates found in registry.")

    pptx_bytes = _potx_to_pptx_bytes(_SOURCE_POTX)

    for entry in entries:
        target = _REPO_ROOT / entry["path"]
        create_template(target, pptx_bytes)

    print(f"Done. {len(entries)} template(s) created.")
