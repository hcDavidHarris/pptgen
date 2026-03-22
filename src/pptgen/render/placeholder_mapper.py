"""Placeholder mapper.

Locates shapes on a slide by their name attribute.

Templates use UPPERCASE_SNAKE_CASE shape names (e.g. TITLE, BULLETS,
METRIC_1_VALUE) as defined in the Template Authoring Standard.  Both
genuine PowerPoint placeholder shapes and regular text-box shapes are
resolved the same way — by matching shape.name.

This design means templates can use named placeholders or named text boxes
interchangeably, as long as the name matches the canonical contract.
"""

from __future__ import annotations

from pptx.shapes.base import BaseShape

from ..errors import TemplateCompatibilityError


def find_placeholder(
    slide,
    placeholder_name: str,
    required: bool = True,
) -> BaseShape | None:
    """Return the shape on *slide* whose name equals *placeholder_name*.

    Args:
        slide:            A python-pptx Slide object.
        placeholder_name: UPPERCASE_SNAKE_CASE name to search for.
        required:         When True (default) raise TemplateCompatibilityError
                          if the shape is not found.  When False, return None.

    Raises:
        TemplateCompatibilityError: if required=True and the shape is absent.
    """
    for shape in slide.shapes:
        if shape.name == placeholder_name:
            return shape

    if not required:
        return None

    available = sorted(s.name for s in slide.shapes)
    raise TemplateCompatibilityError(
        f"Placeholder '{placeholder_name}' not found on slide. "
        f"Available shapes: {available}"
    )


def set_text(slide, placeholder_name: str, text: str) -> None:
    """Set the text of a named shape to *text*, replacing all existing content.

    Uses find_placeholder(required=True) so missing placeholders raise
    TemplateCompatibilityError immediately.
    """
    shape = find_placeholder(slide, placeholder_name)
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = text


def set_bullets(slide, placeholder_name: str, bullets: list[str]) -> None:
    """Populate a named shape with one paragraph per bullet item.

    The first bullet reuses the paragraph that text_frame.clear() leaves
    behind; subsequent bullets are added via add_paragraph().
    """
    shape = find_placeholder(slide, placeholder_name)
    tf = shape.text_frame
    tf.clear()
    for i, bullet_text in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet_text
        else:
            p = tf.add_paragraph()
            p.text = bullet_text
