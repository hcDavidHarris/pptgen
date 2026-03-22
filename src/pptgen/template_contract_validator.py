"""Template contract validator.

Verifies that a PowerPoint template file (.pptx) satisfies the pptgen
placeholder contract for all registered slide types.

The contract requires:
  - Each slide type's layout (e.g. "Bullets Layout") must exist in the template.
  - Each placeholder name (e.g. "TITLE", "BULLETS") must exist as a shape
    name in its respective layout.

This validator is intentionally independent from the deck validator.
It operates on the template file itself and is useful as:
  - a pre-flight check before running pptgen build
  - a CI gate when templates change
  - a diagnostic tool for template authors

Usage::

    from pathlib import Path
    from pptgen.template_contract_validator import validate_template_contract

    result = validate_template_contract(Path("templates/ops_review_v1/template.pptx"))
    if not result.valid:
        for e in result.errors:
            print(f"  ERROR: {e}")
    for w in result.warnings:
        print(f"  WARNING: {w}")
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from .errors import TemplateLoadError
from .slide_registry import SLIDE_TYPE_REGISTRY


@dataclass
class ContractValidationResult:
    """Outcome of a template contract validation pass.

    Attributes:
        valid:    True when no errors were found.
        errors:   Missing layouts or placeholders that will cause render failures.
        warnings: Non-blocking issues found in the template.
    """

    valid: bool
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    def summary(self) -> str:
        """One-line result string."""
        return "PASS" if self.valid else "FAIL"


def validate_template_contract(template_path: Path) -> ContractValidationResult:
    """Validate that *template_path* satisfies the pptgen placeholder contract.

    Checks every slide type registered in SLIDE_TYPE_REGISTRY:
      1. The layout named by ``spec.layout_name`` must exist in the template.
      2. Every placeholder in ``spec.placeholders`` must exist as a shape
         name in that layout.

    Args:
        template_path: Path to the .pptx template file.

    Returns:
        ContractValidationResult with valid=True (PASS) or valid=False (FAIL).

    Raises:
        TemplateLoadError: if the file cannot be opened by python-pptx.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise TemplateLoadError(
            "python-pptx is required for template contract validation"
        ) from exc

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        raise TemplateLoadError(
            f"Cannot open template file '{template_path}': {exc}"
        ) from exc

    # Build a map of layout_name → set of shape names in that layout
    layout_shape_names: dict[str, set[str]] = {}
    for layout in prs.slide_layouts:
        shape_names = {shape.name for shape in layout.shapes}
        layout_shape_names[layout.name] = shape_names

    errors: list[str] = []
    warnings: list[str] = []

    for type_name, spec in SLIDE_TYPE_REGISTRY.items():
        layout_name = spec.layout_name

        # Check layout exists
        if layout_name not in layout_shape_names:
            errors.append(
                f"slide type '{type_name}': layout '{layout_name}' not found in template"
            )
            continue  # Can't check placeholders without the layout

        shapes = layout_shape_names[layout_name]

        # Check each required placeholder exists
        for placeholder_name in spec.placeholders:
            if placeholder_name not in shapes:
                errors.append(
                    f"slide type '{type_name}': layout '{layout_name}': "
                    f"placeholder '{placeholder_name}' not found — "
                    f"available shapes: {sorted(shapes)}"
                )

    return ContractValidationResult(
        valid=len(errors) == 0,
        errors=errors,
        warnings=warnings,
    )
