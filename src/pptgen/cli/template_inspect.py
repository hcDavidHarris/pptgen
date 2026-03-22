"""pptgen template inspect command.

Displays the placeholder contract for a registered template — either from the
slide registry (contract) or by inspecting the actual .pptx file (live shapes).

Usage:
    pptgen template inspect --template ops_review_v1
    pptgen template inspect --template ops_review_v1 --live
"""

from __future__ import annotations

from pathlib import Path

import typer

from ..errors import PptgenError, TemplateLoadError
from ..registry.registry import TemplateRegistry
from ..slide_registry import SLIDE_TYPE_REGISTRY

template_app = typer.Typer(help="Template inspection and management commands.")

_REGISTRY_PATH = Path(__file__).parent.parent.parent.parent / "templates" / "registry.yaml"


@template_app.command(name="inspect")
def inspect(
    template: str = typer.Option(
        ...,
        "--template",
        "-t",
        help="Template ID to inspect. Run 'pptgen list-templates' for available IDs.",
    ),
    live: bool = typer.Option(
        False,
        "--live",
        help="Inspect the actual .pptx file rather than the registry contract.",
    ),
) -> None:
    """Display placeholder contract for a registered template.

    Without --live: shows the expected contract from the slide registry.
    With --live: opens the .pptx file and reports actual shape names per layout.
    """
    try:
        registry = TemplateRegistry.from_file(_REGISTRY_PATH)
    except PptgenError as exc:
        typer.echo(f"Error: {exc}", err=True)
        raise typer.Exit(code=1)

    entry = registry.get(template)
    if entry is None:
        typer.echo(f"Template '{template}' is not registered.", err=True)
        typer.echo("Run 'pptgen list-templates' to see registered templates.")
        raise typer.Exit(code=1)

    typer.echo(f"\nTemplate: {entry.template_id}  v{entry.version}  [{entry.status}]")
    typer.echo(f"Owner:    {entry.owner}")
    typer.echo(f"Path:     {entry.path}")
    typer.echo(f"Supported slide types: {', '.join(entry.supported_slide_types)}")

    if not live:
        typer.echo("\nPlaceholder Contract (from slide registry):")
        typer.echo("-" * 60)
        for type_name in entry.supported_slide_types:
            spec = SLIDE_TYPE_REGISTRY.get(type_name)
            if spec is None:
                continue
            typer.echo(f"\n  {type_name} -> {spec.layout_name}")
            for ph in spec.placeholders:
                typer.echo(f"    {ph}")
        return

    # Live inspection of the .pptx file
    template_path = _REGISTRY_PATH.parent / entry.path
    if not template_path.exists():
        typer.echo(
            f"Template file not found: {template_path}\n"
            "The file may not have been committed yet. "
            "Use the registry contract view (omit --live) instead.",
            err=True,
        )
        raise typer.Exit(code=1)

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        typer.echo("python-pptx is required for live template inspection.", err=True)
        raise typer.Exit(code=1)

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        raise TemplateLoadError(f"Cannot open template: {exc}") from exc

    typer.echo("\nLive Layout Inspection:")
    typer.echo("─" * 60)
    for layout in prs.slide_layouts:
        shape_names = sorted(s.name for s in layout.shapes)
        typer.echo(f"\n  Layout: {layout.name}")
        for name in shape_names:
            typer.echo(f"    {name}")
