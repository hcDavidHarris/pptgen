"""Unit tests for the slide renderer functions.

Each renderer is tested in isolation using a real python-pptx slide built
from the conformant template.  This avoids mocking the pptx layer while
keeping tests fast (no file I/O beyond loading the shared template fixture).
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from pptgen.models.slides import (
    BulletsSlide,
    MetricSummarySlide,
    SectionSlide,
    TitleSlide,
    TwoColumnSlide,
)
from pptgen.render.deck_renderer import (
    SLIDE_TYPE_TO_LAYOUT,
    _rename_slide_placeholders,
)
from pptgen.render.slide_renderers import (
    render_bullets_slide,
    render_metric_summary_slide,
    render_section_slide,
    render_title_slide,
    render_two_column_slide,
    SLIDE_RENDERERS,
)
from pptgen.render.template_inspector import inspect_template


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(scope="module")
def template_prs() -> Presentation:
    """Load the executive_brief_v1 template once per module."""
    path = Path(__file__).parent.parent.parent / "templates" / "executive_brief_v1" / "template.pptx"
    return Presentation(str(path))


@pytest.fixture(scope="module")
def inspection(template_prs):
    return inspect_template(template_prs)


def _make_slide(template_prs, inspection, slide_type: str):
    """Add a slide of *slide_type* from the template and rename its placeholders."""
    layout_name = SLIDE_TYPE_TO_LAYOUT[slide_type]
    layout = inspection.get_layout(layout_name)
    slide = template_prs.slides.add_slide(layout)
    _rename_slide_placeholders(slide, slide_type)
    return slide


def _text(slide, name: str) -> str:
    for shape in slide.shapes:
        if shape.name == name:
            return shape.text_frame.text
    raise KeyError(name)


# ---------------------------------------------------------------------------
# render_title_slide
# ---------------------------------------------------------------------------

class TestRenderTitleSlide:
    def test_sets_title(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "title")
        model = TitleSlide(type="title", title="My Title", subtitle="My Subtitle")
        render_title_slide(model, slide)
        assert _text(slide, "TITLE") == "My Title"

    def test_sets_subtitle(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "title")
        model = TitleSlide(type="title", title="T", subtitle="S")
        render_title_slide(model, slide)
        assert _text(slide, "SUBTITLE") == "S"

    def test_subtitle_unicode(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "title")
        model = TitleSlide(type="title", title="T", subtitle="Q2 — FY2025")
        render_title_slide(model, slide)
        assert _text(slide, "SUBTITLE") == "Q2 — FY2025"


# ---------------------------------------------------------------------------
# render_section_slide
# ---------------------------------------------------------------------------

class TestRenderSectionSlide:
    def test_sets_section_title(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "section")
        model = SectionSlide(type="section", section_title="Intro")
        render_section_slide(model, slide)
        assert _text(slide, "SECTION_TITLE") == "Intro"

    def test_section_subtitle_written_when_present(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "section")
        model = SectionSlide(type="section", section_title="S", section_subtitle="desc")
        render_section_slide(model, slide)
        assert _text(slide, "SECTION_SUBTITLE") == "desc"

    def test_section_subtitle_empty_when_absent(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "section")
        model = SectionSlide(type="section", section_title="S")
        render_section_slide(model, slide)
        assert _text(slide, "SECTION_SUBTITLE") == ""


# ---------------------------------------------------------------------------
# render_bullets_slide
# ---------------------------------------------------------------------------

class TestRenderBulletsSlide:
    def test_sets_title(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "bullets")
        model = BulletsSlide(type="bullets", title="Points", bullets=["A"])
        render_bullets_slide(model, slide)
        assert _text(slide, "TITLE") == "Points"

    def test_single_bullet(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "bullets")
        model = BulletsSlide(type="bullets", title="T", bullets=["Only one"])
        render_bullets_slide(model, slide)
        assert _text(slide, "BULLETS") == "Only one"

    def test_multiple_bullets_joined_by_newline(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "bullets")
        model = BulletsSlide(type="bullets", title="T", bullets=["A", "B", "C"])
        render_bullets_slide(model, slide)
        assert _text(slide, "BULLETS") == "A\nB\nC"


# ---------------------------------------------------------------------------
# render_two_column_slide
# ---------------------------------------------------------------------------

class TestRenderTwoColumnSlide:
    def test_sets_title_and_columns(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "two_column")
        model = TwoColumnSlide(
            type="two_column", title="Compare",
            left_content=["L1", "L2"], right_content=["R1"],
        )
        render_two_column_slide(model, slide)
        assert _text(slide, "TITLE") == "Compare"
        assert _text(slide, "LEFT_CONTENT") == "L1\nL2"
        assert _text(slide, "RIGHT_CONTENT") == "R1"


# ---------------------------------------------------------------------------
# render_metric_summary_slide
# ---------------------------------------------------------------------------

class TestRenderMetricSummarySlide:
    def _make_model(self, metrics):
        from pptgen.models.slides import MetricItem
        return MetricSummarySlide(
            type="metric_summary",
            title="KPIs",
            metrics=[MetricItem(**m) for m in metrics],
        )

    def test_title_set(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "metric_summary")
        model = self._make_model([{"label": "A", "value": "1"}])
        render_metric_summary_slide(model, slide)
        assert _text(slide, "TITLE") == "KPIs"

    def test_first_metric_populated(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "metric_summary")
        model = self._make_model([{"label": "Uptime", "value": "99.9", "unit": "%"}])
        render_metric_summary_slide(model, slide)
        assert _text(slide, "METRIC_1_LABEL") == "Uptime"
        assert _text(slide, "METRIC_1_VALUE") == "99.9%"

    def test_unit_concatenated_directly(self, template_prs, inspection):
        """No separator is added — any space must be in the unit string itself."""
        slide = _make_slide(template_prs, inspection, "metric_summary")
        model = self._make_model([{"label": "Latency", "value": "42", "unit": " ms"}])
        render_metric_summary_slide(model, slide)
        assert _text(slide, "METRIC_1_VALUE") == "42 ms"

    def test_unit_none_leaves_value_unchanged(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "metric_summary")
        model = self._make_model([{"label": "L", "value": "5"}])
        render_metric_summary_slide(model, slide)
        assert _text(slide, "METRIC_1_VALUE") == "5"

    def test_unused_positions_cleared(self, template_prs, inspection):
        """Positions beyond the model's metric count must be empty strings."""
        slide = _make_slide(template_prs, inspection, "metric_summary")
        model = self._make_model([{"label": "Only", "value": "1"}])
        render_metric_summary_slide(model, slide)
        for pos in (2, 3, 4):
            assert _text(slide, f"METRIC_{pos}_LABEL") == ""
            assert _text(slide, f"METRIC_{pos}_VALUE") == ""

    def test_four_metrics_all_populated(self, template_prs, inspection):
        slide = _make_slide(template_prs, inspection, "metric_summary")
        model = self._make_model([
            {"label": f"L{i}", "value": str(i)} for i in range(1, 5)
        ])
        render_metric_summary_slide(model, slide)
        for pos in (1, 2, 3, 4):
            assert _text(slide, f"METRIC_{pos}_LABEL") == f"L{pos}"
            assert _text(slide, f"METRIC_{pos}_VALUE") == str(pos)


# ---------------------------------------------------------------------------
# SLIDE_RENDERERS registry
# ---------------------------------------------------------------------------

class TestSlideRenderersRegistry:
    def test_all_types_registered(self):
        expected = {"title", "section", "bullets", "two_column", "metric_summary", "image_caption"}
        assert expected.issubset(SLIDE_RENDERERS.keys())

    def test_values_are_callable(self):
        for fn in SLIDE_RENDERERS.values():
            assert callable(fn)
