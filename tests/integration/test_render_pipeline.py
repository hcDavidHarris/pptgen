"""Integration tests for the full render pipeline.

Tests run the complete path from YAML → DeckFile → .pptx, verifying that
the output presentation has the expected slide count and placeholder content.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from pptgen.loaders.yaml_loader import load_deck
from pptgen.registry.registry import TemplateRegistry
from pptgen.render import render_deck
from pptgen.validators.deck_validator import validate_deck


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_PROJECT_ROOT = Path(__file__).parent.parent.parent
_REGISTRY_PATH = _PROJECT_ROOT / "templates" / "registry.yaml"
_EXAMPLES_DIR = _PROJECT_ROOT / "examples"


def _slide_text(prs: Presentation, slide_idx: int, shape_name: str) -> str:
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape.text_frame.text
    raise KeyError(f"Shape '{shape_name}' not found on slide {slide_idx}")


@pytest.fixture(scope="module")
def registry() -> TemplateRegistry:
    return TemplateRegistry.from_file(_REGISTRY_PATH)


@pytest.fixture(scope="module")
def executive_update_prs(tmp_path_factory, registry) -> Presentation:
    """Render executive_update.yaml to a temp .pptx and return the opened result."""
    deck, raw = load_deck(_EXAMPLES_DIR / "executive_update.yaml")
    result = validate_deck(deck, registry, raw)
    assert result.valid, f"Validation failed: {result.errors}"

    entry = registry.get(deck.deck.template)
    template_path = _PROJECT_ROOT / entry.path
    output_path = tmp_path_factory.mktemp("render") / "executive_update.pptx"

    render_deck(deck, template_path, output_path)
    return Presentation(str(output_path))


# ---------------------------------------------------------------------------
# executive_update.yaml  (6 slides: title, section, bullets, section, bullets, metric_summary)
# ---------------------------------------------------------------------------

class TestExecutiveUpdateRender:
    def test_correct_slide_count(self, executive_update_prs):
        assert len(executive_update_prs.slides) == 6

    def test_title_slide_title(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 0, "TITLE") == "Executive Update"

    def test_title_slide_subtitle(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 0, "SUBTITLE") == "Quarterly Progress Review"

    def test_first_section_title(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 1, "SECTION_TITLE") == "Business Context"

    def test_first_section_subtitle(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 1, "SECTION_SUBTITLE") == (
            "Current priorities and operating environment"
        )

    def test_bullets_slide_title(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 2, "TITLE") == "Key Highlights"

    def test_bullets_slide_content(self, executive_update_prs):
        text = _slide_text(executive_update_prs, 2, "BULLETS")
        assert "Platform stability improved" in text
        assert "Reporting automation" in text
        assert "Adoption increased" in text

    def test_section_without_subtitle_is_empty(self, executive_update_prs):
        """Slide 4 (index 3) has no section_subtitle — placeholder should be empty."""
        assert _slide_text(executive_update_prs, 3, "SECTION_TITLE") == "Strategic Priorities"
        assert _slide_text(executive_update_prs, 3, "SECTION_SUBTITLE") == ""

    def test_metric_summary_title(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 5, "TITLE") == "KPI Snapshot"

    def test_metric_summary_first_metric(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 5, "METRIC_1_LABEL") == "Adoption Growth"
        assert _slide_text(executive_update_prs, 5, "METRIC_1_VALUE") == "27%"

    def test_metric_summary_third_metric(self, executive_update_prs):
        assert _slide_text(executive_update_prs, 5, "METRIC_3_LABEL") == "Delivery Reliability"
        assert _slide_text(executive_update_prs, 5, "METRIC_3_VALUE") == "99.2%"

    def test_metric_summary_unused_position_empty(self, executive_update_prs):
        """Position 4 has no metric — should be cleared to empty string."""
        assert _slide_text(executive_update_prs, 5, "METRIC_4_LABEL") == ""
        assert _slide_text(executive_update_prs, 5, "METRIC_4_VALUE") == ""


# ---------------------------------------------------------------------------
# visible=false slides are skipped
# ---------------------------------------------------------------------------

class TestVisibleFlag:
    def test_hidden_slide_excluded_from_output(self, tmp_path, registry):
        """A deck with one visible=false slide produces one fewer output slide."""
        import yaml
        deck_data = {
            "deck": {"title": "T", "template": "executive_brief_v1", "author": "A"},
            "slides": [
                {"type": "title", "title": "Visible", "subtitle": "yes"},
                {"type": "title", "title": "Hidden", "subtitle": "no", "visible": False},
            ],
        }
        yaml_path = tmp_path / "deck.yaml"
        yaml_path.write_text(yaml.dump(deck_data))

        deck, raw = load_deck(yaml_path)
        entry = registry.get(deck.deck.template)
        template_path = _PROJECT_ROOT / entry.path
        output_path = tmp_path / "out.pptx"

        render_deck(deck, template_path, output_path)
        prs = Presentation(str(output_path))
        assert len(prs.slides) == 1
        assert _slide_text(prs, 0, "TITLE") == "Visible"
