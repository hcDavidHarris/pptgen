"""Integration tests for image_caption slide rendering.

Verifies the full render path: YAML deck → DeckFile → render_deck() → .pptx
for slides of type image_caption.  The image_caption renderer was added in
the Phase 3 Closeout cycle; these tests guard against regression.
"""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation

from pptgen.loaders.yaml_loader import load_deck
from pptgen.registry.registry import TemplateRegistry
from pptgen.render import render_deck
from pptgen.render.slide_renderers import SLIDE_RENDERERS


_PROJECT_ROOT = Path(__file__).parent.parent
_REGISTRY_PATH = _PROJECT_ROOT / "templates" / "registry.yaml"
_TEMPLATE_PATH = _PROJECT_ROOT / "templates" / "ops_review_v1" / "template.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _shape_text(prs: Presentation, slide_idx: int, shape_name: str) -> str:
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape.text_frame.text
    raise KeyError(f"Shape '{shape_name}' not found on slide {slide_idx}")


def _make_image_caption_deck_yaml(image_path: str = "assets/placeholder.png") -> dict:
    return {
        "deck": {
            "title": "Image Caption Test Deck",
            "template": "ops_review_v1",
            "author": "Test Runner",
            "version": "1.0",
        },
        "slides": [
            {
                "type": "title",
                "title": "Image Caption Test",
                "subtitle": "Integration test",
            },
            {
                "type": "image_caption",
                "id": "arch_slide",
                "title": "Architecture Overview",
                "image_path": image_path,
                "caption": "System architecture diagram showing all core components",
            },
        ],
    }


# ---------------------------------------------------------------------------
# Unit: renderer registry
# ---------------------------------------------------------------------------

class TestImageCaptionRendererRegistered:
    def test_image_caption_in_slide_renderers(self):
        assert "image_caption" in SLIDE_RENDERERS

    def test_renderer_is_callable(self):
        assert callable(SLIDE_RENDERERS["image_caption"])


# ---------------------------------------------------------------------------
# Integration: full render pipeline
# ---------------------------------------------------------------------------

class TestImageCaptionRenderPipeline:
    @pytest.fixture(scope="class")
    def rendered_prs(self, tmp_path_factory) -> Presentation:
        """Render a deck containing an image_caption slide; return opened result."""
        deck_data = _make_image_caption_deck_yaml()
        yaml_path = tmp_path_factory.mktemp("image_caption") / "deck.yaml"
        yaml_path.write_text(yaml.dump(deck_data), encoding="utf-8")

        registry = TemplateRegistry.from_file(_REGISTRY_PATH)
        deck, raw = load_deck(yaml_path)

        output_path = yaml_path.parent / "out.pptx"
        render_deck(deck, _TEMPLATE_PATH, output_path)
        return Presentation(str(output_path))

    def test_output_file_created(self, tmp_path):
        deck_data = _make_image_caption_deck_yaml()
        yaml_path = tmp_path / "deck.yaml"
        yaml_path.write_text(yaml.dump(deck_data), encoding="utf-8")

        output_path = tmp_path / "out.pptx"
        deck, _ = load_deck(yaml_path)
        render_deck(deck, _TEMPLATE_PATH, output_path)

        assert output_path.exists()

    def test_correct_slide_count(self, rendered_prs):
        assert len(rendered_prs.slides) == 2

    def test_title_placeholder_written(self, rendered_prs):
        assert _shape_text(rendered_prs, 1, "TITLE") == "Architecture Overview"

    def test_caption_placeholder_written(self, rendered_prs):
        text = _shape_text(rendered_prs, 1, "CAPTION")
        assert text == "System architecture diagram showing all core components"


# ---------------------------------------------------------------------------
# Integration: image_caption with missing image (graceful degradation)
# ---------------------------------------------------------------------------

class TestImageCaptionMissingImage:
    def test_render_succeeds_when_image_path_does_not_exist(self, tmp_path):
        """Renderer must not raise when image_path points to a non-existent file.

        The image placeholder is left empty; text placeholders are still written.
        """
        deck_data = _make_image_caption_deck_yaml(
            image_path="assets/does_not_exist.png"
        )
        yaml_path = tmp_path / "deck.yaml"
        yaml_path.write_text(yaml.dump(deck_data), encoding="utf-8")

        output_path = tmp_path / "out.pptx"
        deck, _ = load_deck(yaml_path)
        # Must not raise KeyError or FileNotFoundError
        render_deck(deck, _TEMPLATE_PATH, output_path)

        assert output_path.exists()
        prs = Presentation(str(output_path))
        assert _shape_text(prs, 1, "TITLE") == "Architecture Overview"
